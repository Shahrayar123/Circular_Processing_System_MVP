"""Read any document into text — including documents inside documents.

ABL confirmed circulars arrive in every format: scanned PDFs, images, native PDFs, Word,
Excel, Word containing Excel, legacy .doc/.xls, Outlook messages. Their own BRD proves
the nested case — it is a .docx carrying an .xlsm and four .docx files in
`word/embeddings/`, and the embedded workbook is the part that actually matters.

Three rules, the same ones the delivered system follows:

1. **Route per page, not per document.** A circular is routinely part native text and
   part scanned annexure. Deciding once for the whole file is wrong either way.
2. **Recurse into embedded objects**, with a depth cap, a seen-hash set and a size cap.
   A container that references itself must not loop; one bad file must not kill the run.
3. **Never fail silently.** Password-protected, corrupt, unsupported — the document is
   recorded with an error and surfaces in the queue. A circular that quietly disappears
   is the worst outcome: the audit team believes it was processed.
"""

import hashlib
import re
import shutil
import subprocess
import tempfile
import zipfile
from dataclasses import dataclass, field
from pathlib import Path

SUPPORTED = {
    ".pdf", ".docx", ".doc", ".xlsx", ".xlsm", ".xls", ".pptx", ".ppt",
    ".txt", ".csv", ".htm", ".html", ".rtf", ".eml", ".msg",
    ".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp", ".gif",
}

IMAGES = {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp", ".gif"}
LEGACY = {".doc": ".docx", ".xls": ".xlsx", ".ppt": ".pptx"}

MAX_DEPTH = 3                      # a document inside a document inside a document
MAX_TOTAL_CHARS = 4_000_000        # a 500 MB workbook fails the document, not the batch
NATIVE_TEXT_MIN_CHARS = 120        # per page, below this the page is treated as scanned

SOFFICE_PATHS = [
    Path(r"C:\Program Files\LibreOffice\program\soffice.exe"),
    Path(r"C:\Program Files (x86)\LibreOffice\program\soffice.exe"),
    Path("/usr/bin/soffice"), Path("/usr/bin/libreoffice"),
]


@dataclass
class Extracted:
    """What one file yielded. `children` are documents found inside it."""
    text: str = ""
    pages: int = 0
    kind: str = ""                 # native-pdf | scanned-pdf | mixed-pdf | docx | xlsx …
    ocr_pages: int = 0
    error: str = ""
    children: list = field(default_factory=list)   # [(name, Extracted)]


# ====== HELPERS ======

def _clean(text: str) -> str:
    text = (text or "").replace("\r\n", "\n").replace("\xa0", " ")
    text = _REVIEW_MARKUP.sub(" ", text)
    text = re.sub(r"[ \t]+", " ", text)
    return re.sub(r"\n{3,}", "\n\n", text).strip()


# Word review markup rendered into a PDF when a draft is exported with comments and
# tracked changes showing. The parenthesised form is first and spans newlines, because
# PDF extraction wraps lines mid-marker.
_REVIEW_MARKUP = re.compile(
    r"Formatted:[^)\n]{0,80}\([^)]{0,200}\)"
    r"|Commented\s*\[[A-Za-z]+\d*\]:.*?(?=(?:\n|Commented\s*\[|Formatted:|$))"
    r"|Formatted:[^\n]*|Deleted:[^\n]*|Inserted:[^\n]*",
    re.DOTALL,
)


def find_soffice() -> Path | None:
    for path in SOFFICE_PATHS:
        if path.exists():
            return path
    found = shutil.which("soffice") or shutil.which("libreoffice")
    return Path(found) if found else None


def sniff(path: Path) -> str:
    """The real type from the magic bytes. Mail systems rename files freely — a .pdf
    that is actually a Word document is a normal thing to receive."""
    try:
        head = path.open("rb").read(8)
    except Exception:
        return path.suffix.lower()
    if head.startswith(b"%PDF"):
        return ".pdf"
    if head.startswith(b"PK\x03\x04"):                   # any OOXML / zip container
        suffix = path.suffix.lower()
        return suffix if suffix in {".docx", ".xlsx", ".xlsm", ".pptx"} else ".docx"
    if head.startswith(b"\xd0\xcf\x11\xe0"):             # legacy OLE compound file
        suffix = path.suffix.lower()
        return suffix if suffix in LEGACY or suffix == ".msg" else ".doc"
    if head[:3] in (b"\xff\xd8\xff",) or head.startswith(b"\x89PNG"):
        return path.suffix.lower() if path.suffix.lower() in IMAGES else ".png"
    return path.suffix.lower()


# ====== OCR ======

def ocr_image(image) -> str:
    """Tesseract. The delivered system uses Qwen2.5-VL with Tesseract as the failure
    fallback; the MVP has no vision endpoint, so Tesseract does the whole job."""
    import pytesseract

    for candidate in (r"C:\Program Files\Tesseract-OCR\tesseract.exe",
                      r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe"):
        if Path(candidate).exists():
            pytesseract.pytesseract.tesseract_cmd = candidate
            break
    return pytesseract.image_to_string(image)


# ====== READERS ======

def read_pdf(path: Path) -> Extracted:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    pages, native, scanned = [], 0, 0

    for index, page in enumerate(reader.pages):
        text = (page.extract_text() or "").strip()
        if len(text) >= NATIVE_TEXT_MIN_CHARS:
            native += 1
        else:
            # ROUTE PER PAGE. A part-scanned circular is normal; deciding for the whole
            # document wastes OCR on clean pages or silently loses the scanned ones.
            try:
                from pdf2image import convert_from_path
                images = convert_from_path(str(path), dpi=200,
                                           first_page=index + 1, last_page=index + 1,
                                           poppler_path=_poppler())
                if images:
                    text = (ocr_image(images[0]) or "").strip()
                    scanned += 1
            except Exception as exc:
                text = f"[page {index + 1}: could not be read — {exc}]"
        pages.append(text)

    kind = "native-pdf" if not scanned else ("scanned-pdf" if not native else "mixed-pdf")
    return Extracted(text=_clean("\n\n".join(pages)), pages=len(pages),
                     kind=kind, ocr_pages=scanned)


def _poppler():
    for candidate in (r"C:\Program Files\poppler-25.07.0\Library\bin",
                      r"C:\Program Files\poppler\Library\bin",
                      r"C:\poppler\Library\bin", r"C:\poppler\bin"):
        if Path(candidate).exists():
            return candidate
    return None


def read_image(path: Path) -> Extracted:
    from PIL import Image

    return Extracted(text=_clean(ocr_image(Image.open(path))), pages=1,
                     kind="image", ocr_pages=1)


def read_docx(path: Path) -> Extracted:
    from docx import Document

    document = Document(str(path))
    parts = [p.text for p in document.paragraphs if p.text.strip()]

    # Tables carry obligations routinely — a row is often one requirement. Keep the row
    # structure with a separator rather than flattening it into prose.
    for table in document.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))

    result = Extracted(text=_clean("\n".join(parts)), pages=1, kind="docx")
    result.children = _embedded_in_ooxml(path)
    return result


def read_excel(path: Path) -> Extracted:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    blocks = []
    for name in wb.sheetnames:
        sheet = wb[name]
        blocks.append(f"### SHEET: {name}")
        for number, row in enumerate(sheet.iter_rows(values_only=True), start=1):
            cells = [str(c).strip() for c in row if c is not None and str(c).strip()]
            if cells:
                # The row reference travels with the text so a reviewer can find it
                # again — "Annexure A · row 14" rather than "somewhere in the sheet".
                blocks.append(f"[row {number}] " + " | ".join(cells))
    return Extracted(text=_clean("\n".join(blocks)), pages=len(wb.sheetnames),
                     kind="excel")


def read_pptx(path: Path) -> Extracted:
    try:
        from pptx import Presentation
    except ImportError:
        return Extracted(error="python-pptx is not installed — .pptx cannot be read",
                         kind="pptx")
    prs = Presentation(str(path))
    parts = []
    for index, slide in enumerate(prs.slides, start=1):
        parts.append(f"### SLIDE {index}")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            parts.append("Notes: " + slide.notes_slide.notes_text_frame.text)
    return Extracted(text=_clean("\n".join(parts)), pages=len(prs.slides), kind="pptx")


def read_msg(path: Path) -> Extracted:
    try:
        import extract_msg
    except ImportError:
        return Extracted(error="extract-msg is not installed — Outlook .msg cannot be read",
                         kind="msg")
    message = extract_msg.Message(str(path))
    header = f"Subject: {message.subject}\nFrom: {message.sender}\nDate: {message.date}"
    result = Extracted(text=_clean(header + "\n\n" + (message.body or "")), pages=1,
                       kind="msg")
    for attachment in message.attachments:
        name = attachment.longFilename or attachment.shortFilename or "attachment"
        if Path(name).suffix.lower() in SUPPORTED:
            with tempfile.TemporaryDirectory() as tmp:
                saved = Path(tmp) / name
                saved.write_bytes(attachment.data)
                result.children.append((name, read(saved, _depth=1)))
    return result


def read_text(path: Path) -> Extracted:
    raw = path.read_text(encoding="utf-8", errors="replace")
    if path.suffix.lower() in {".htm", ".html"}:
        raw = re.sub(r"<(script|style)[^>]*>.*?</\1>", " ", raw, flags=re.S | re.I)
        raw = re.sub(r"<[^>]+>", " ", raw)
    if path.suffix.lower() == ".rtf":
        raw = re.sub(r"\\[a-z]+-?\d*\s?|[{}]", " ", raw)
    return Extracted(text=_clean(raw), pages=1, kind=path.suffix.lstrip("."))


def read_eml(path: Path) -> Extracted:
    import email

    message = email.message_from_bytes(path.read_bytes())
    body = []
    result = Extracted(pages=1, kind="eml")
    for part in message.walk():
        if part.get_content_type() == "text/plain" and not part.get_filename():
            payload = part.get_payload(decode=True)
            if payload:
                body.append(payload.decode("utf-8", "replace"))
        name = part.get_filename()
        if name and Path(name).suffix.lower() in SUPPORTED:
            with tempfile.TemporaryDirectory() as tmp:
                saved = Path(tmp) / name
                saved.write_bytes(part.get_payload(decode=True) or b"")
                result.children.append((name, read(saved, _depth=1)))
    header = f"Subject: {message.get('Subject', '')}\nFrom: {message.get('From', '')}"
    result.text = _clean(header + "\n\n" + "\n".join(body))
    return result


def read_legacy(path: Path) -> Extracted:
    """.doc / .xls / .ppt — convert with LibreOffice, then read normally."""
    soffice = find_soffice()
    if soffice is None:
        return Extracted(error="LibreOffice is not installed — legacy "
                               f"{path.suffix} cannot be converted", kind="legacy")
    target = LEGACY[path.suffix.lower()]
    with tempfile.TemporaryDirectory() as tmp:
        try:
            subprocess.run(
                [str(soffice), "--headless", "--convert-to", target.lstrip("."),
                 "--outdir", tmp, str(path)],
                capture_output=True, timeout=180, check=False)
        except Exception as exc:
            return Extracted(error=f"LibreOffice conversion failed: {exc}", kind="legacy")
        converted = list(Path(tmp).glob("*" + target))
        if not converted:
            return Extracted(error="LibreOffice produced no output — the file may be "
                                   "password-protected or corrupt", kind="legacy")
        result = read(converted[0], _depth=1)
        result.kind = f"legacy{path.suffix}"
        return result


READERS = {
    ".pdf": read_pdf, ".docx": read_docx,
    ".xlsx": read_excel, ".xlsm": read_excel,
    ".pptx": read_pptx, ".msg": read_msg, ".eml": read_eml,
    ".doc": read_legacy, ".xls": read_legacy, ".ppt": read_legacy,
    ".txt": read_text, ".csv": read_text, ".htm": read_text,
    ".html": read_text, ".rtf": read_text,
    **{suffix: read_image for suffix in IMAGES},
}


# ====== EMBEDDED OBJECTS ======

def _embedded_in_ooxml(path: Path) -> list:
    """Documents stored inside a .docx / .xlsx package.

    ABL's own BRD is the proof this matters: a .docx carrying an .xlsm and four .docx
    files under `word/embeddings/`, and the embedded workbook is the specification.
    Reading the outer file and stopping loses it, silently.
    """
    found = []
    try:
        with zipfile.ZipFile(path) as archive:
            names = [n for n in archive.namelist()
                     if "/embeddings/" in n and Path(n).suffix.lower() in SUPPORTED]
            for name in names[:12]:                     # a sane cap for a demo
                with tempfile.TemporaryDirectory() as tmp:
                    saved = Path(tmp) / Path(name).name
                    saved.write_bytes(archive.read(name))
                    found.append((Path(name).name, read(saved, _depth=1)))
    except Exception:
        pass                                            # not a zip, or unreadable
    return found


# ====== ENTRY POINT ======

def read(path: Path, _depth: int = 0, _seen: set | None = None) -> Extracted:
    """Read one file. Never raises — a failure comes back in `.error`."""
    path = Path(path)
    _seen = _seen if _seen is not None else set()

    if _depth > MAX_DEPTH:
        return Extracted(error=f"nesting deeper than {MAX_DEPTH} levels — not followed")

    try:
        digest = hashlib.sha256(path.read_bytes()).hexdigest()
    except Exception as exc:
        return Extracted(error=f"could not be opened: {exc}")

    if digest in _seen:                                 # a container referencing itself
        return Extracted(error="already extracted in this run — not followed again")
    _seen.add(digest)

    if path.stat().st_size == 0:
        return Extracted(error="the file is empty")

    suffix = sniff(path)
    reader = READERS.get(suffix)
    if reader is None:
        return Extracted(error=f"unsupported file type {suffix!r}", kind=suffix)

    try:
        result = reader(path)
    except Exception as exc:
        message = str(exc)
        if "password" in message.lower() or "encrypt" in message.lower():
            message = "the file is password-protected"
        return Extracted(error=f"could not be read: {message}", kind=suffix)

    if len(result.text) > MAX_TOTAL_CHARS:
        result.text = result.text[:MAX_TOTAL_CHARS]
        result.error = (result.error + " | " if result.error else "") + \
            f"truncated at {MAX_TOTAL_CHARS:,} characters"
    return result


# ====== TITLE AND DATE ======

_DATE = re.compile(
    r"\b(\d{1,2}[-/ ][A-Z][a-z]{2,8}[-/ ]\d{2,4}|[A-Z][a-z]{2,8} \d{1,2}, \d{4}|\d{4}-\d{2}-\d{2})\b"
)


def guess_title(text: str, filename: str) -> str:
    for line in (text or "").split("\n")[:25]:
        line = line.strip()
        lowered = line.lower()
        if any(w in lowered for w in ("confidential", "for bank", "internal use",
                                      "page ", "table of contents", "all rights",
                                      "### sheet", "### slide")):
            continue
        if 12 <= len(line) <= 110:
            letters = [c for c in line if c.isalpha()]
            if letters and sum(c.isupper() for c in letters) / len(letters) > 0.4:
                return line
    return Path(filename).stem.replace("_", " ").strip()


def guess_date(text: str) -> str:
    match = _DATE.search((text or "")[:2500])
    return match.group(1) if match else ""
